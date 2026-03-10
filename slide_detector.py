"""Slide detection functionality."""

import logging
import re
from dataclasses import dataclass
from typing import Generator, List, Optional
from pptx import Presentation

from config import ISSUE_PATTERNS, ISSUE_PROJECT_RULES, DEFAULT_PROJECT_KEY

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class IssueSlideReference:
    """Coordinates for an issue slide across the PPTX and exported PDF."""

    pptx_slide_number: int
    pdf_page_number: int
    project_key: str


class SlideDetector:
    """Handles detection of issue slides in presentations."""
    
    def __init__(self, patterns: List[str] = None):
        self.patterns = patterns or ISSUE_PATTERNS
    
    def find_issue_slides(self, pptx_path: str) -> Generator[IssueSlideReference, None, None]:
        """Yield issue slides with PPTX numbering and exported PDF page numbering."""
        try:
            prs = Presentation(pptx_path)
            logger.info(f"Processing {len(prs.slides)} slides from {pptx_path}")

            visible_slide_number = 0

            for idx, slide in enumerate(prs.slides, start=1):
                is_hidden = self._is_hidden(slide)
                if not is_hidden:
                    visible_slide_number += 1

                project_key = self._detect_issue_and_project(slide)
                if project_key is not None:  # Found an issue
                    if is_hidden:
                        logger.warning(
                            "Skipping hidden issue slide %s because hidden slides are not exported to PDF",
                            idx,
                        )
                        continue

                    logger.info(
                        "Found issue slide: pptx=%s pdf=%s → project: %s",
                        idx,
                        visible_slide_number,
                        project_key,
                    )
                    yield IssueSlideReference(
                        pptx_slide_number=idx,
                        pdf_page_number=visible_slide_number,
                        project_key=project_key,
                    )
        except Exception as e:
            logger.error(f"Error processing presentation: {e}")
            raise

    def _is_hidden(self, slide) -> bool:
        """Return True when a slide is marked hidden in the PPTX XML."""
        show_attr = slide._element.get("show")
        if show_attr is None:
            return False

        return str(show_attr).strip().lower() in {"0", "false"}
    
    def _detect_issue_and_project(self, slide) -> Optional[str]:
        """Detect issue slides and determine project key based on text patterns."""
        slide_text = self._extract_slide_text(slide)
        
        # First check if it's an issue slide at all
        if not any(re.search(pattern, slide_text) for pattern in self.patterns):
            return None  # Not an issue slide
        
        # It's an issue slide, now determine the project
        for pattern, project_key in ISSUE_PROJECT_RULES.items():
            if re.search(pattern, slide_text):
                logger.debug(f"Project rule matched: '{pattern}' → {project_key}")
                return project_key
        
        # Default project if no specific rules match
        logger.debug(f"No project rule matched, using default: {DEFAULT_PROJECT_KEY}")
        return DEFAULT_PROJECT_KEY
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text from a slide."""
        texts = []
        for shp in slide.shapes:
            if hasattr(shp, "text") and shp.text.strip():
                texts.append(shp.text.strip())
        return "\n".join(texts) 
